import re
import copy
import docx
import pptx
import pathlib
import openpyxl
import functools
import multiprocessing
import deep_translator
import pebble.concurrent

import utilities


def translate_str_vul(text: str) -> str:
    return deep_translator.GoogleTranslator(source="auto", target="en").translate(text)


@pebble.concurrent.process(timeout=30)
def translate_str_sec(text: str) -> str:
    return deep_translator.GoogleTranslator(source="auto", target="en").translate(text)


@functools.cache
def translate_str(safe: bool, text: str) -> str:
    return translate_str_vul(text) if not safe else translate_str_sec(text).result()


def translate_xls(safe: bool, path: pathlib.Path, watch: utilities.Watch) -> None:

    def _translate(container) -> None:
        if isinstance(container.value, str) and re.search("[\u4E00-\u9FFF]", container.value) and container.data_type != "f":
            try:
                translation = translate_str(safe, container.value)
                if isinstance(translation, str): container.value = translation
            except Exception as exception:
                watch.print(f"{exception}")

    workbook = openpyxl.load_workbook(str(path))
    for sheetname, (rows, cols) in utilities.worksheets_dimensions(str(path)).items():
        for row in range(1, rows + 2):
            for col in range(1, cols + 1):
                cell = workbook[sheetname].cell(row, col)
                _translate(cell)
    workbook.save(str(path))


def translate_doc(safe: bool, path: pathlib.Path, watch: utilities.Watch) -> None:

    def _translate(container) -> None:
        if isinstance(container.text, str) and re.search("[\u4E00-\u9FFF]", container.text):
            try:
                translation = translate_str(safe, container.text)
                if isinstance(translation, str): container.text = translation
            except Exception as exception:
                watch.print(f"{exception}")

    document = docx.Document(str(path))
    paragraphs = document.paragraphs
    tables = document.tables
    for section in document.sections:
        paragraphs += section.header.paragraphs + section.footer.paragraphs
        tables += section.header.tables + section.footer.tables
    for paragraph in paragraphs:
        for run in paragraph.runs:
            _translate(run)
    for table in tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        _translate(run)
                for subtable in cell.tables:
                    for subrow in subtable.rows:
                        for subcell in subrow.cells:
                            for subparagraph in subcell.paragraphs:
                                for subrun in subparagraph.runs:
                                    _translate(subrun)
    document.save(str(path))


def translate_ppt(safe: bool, path: pathlib.Path, watch: utilities.Watch) -> None:

    def _translate(container) -> None:
        if isinstance(container.text, str) and re.search("[\u4E00-\u9FFF]", container.text):
            try:
                translation = translate_str(safe, container.text)
                if isinstance(translation, str): container.text = translation
            except Exception as exception:
                watch.print(f"{exception}")

    presentation = pptx.Presentation(str(path))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        _translate(run)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                _translate(run)
    presentation.save(str(path))


def translate_mso(safe: bool, paths: list[tuple[int, pathlib.Path]], watch: utilities.Watch) -> None:

    suffix_to_function = {".xlsx": translate_xls, ".docx": translate_doc, ".pptx": translate_ppt}

    while paths:
        index, path = paths.pop(0)
        watch.current(index, path)
        watch.print("translating")
        function = suffix_to_function[path.suffix]
        function(safe, path, copy.deepcopy(watch))
        watch.print("translated")


def translate(safe: bool, directory_dst: pathlib.Path, processes: int, watch: utilities.Watch) -> None:

    with multiprocessing.Manager() as manager:

        paths_raw = [path for path in directory_dst.rglob("*") if path.suffix in (".xlsx", ".docx", ".pptx")]

        paths = manager.list(enumerate(paths_raw, start=1))
        watch.beginning(len(paths_raw))

        pool = [multiprocessing.Process(target=translate_mso, args=(safe, paths, copy.deepcopy(watch))) for _ in range(processes)]
        for process in pool: process.start()
        for process in pool: process.join()
